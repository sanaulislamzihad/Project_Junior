import pdfplumber
import os

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extracts all text from a PDF file.
    
    Args:
        file_path (str): The absolute path to the PDF file.
        
    Returns:
        str: The extracted text content.
    """
    full_text = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    full_text.append(text)
        
        combined_text = "\n\n".join(full_text)
        
        # Clean up common PDF artifacts and ligatures
        replacements = {
            # Ligatures
            "ﬁ": "fi", "ﬂ": "fl", "ﬀ": "ff", "ﬃ": "ffi", "ﬄ": "ffl",
            # Smart quotes and dashes
            "“": '"', "”": '"', "‘": "'", "’": "'",
            "–": "-", "—": "-",
            # Zero width space
            "\u200b": ""
        }
        for old, new in replacements.items():
            combined_text = combined_text.replace(old, new)

        # Return text with original layout (newlines preserved)
        # We process this normalization later in main.py for matching
        return combined_text
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PowerPoint file."""
    from pptx import Presentation
    full_text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        combined_text = "\n\n".join(full_text)
        return clean_text_artifacts(combined_text)
    except Exception as e:
        print(f"Error extracting text from PPTX {file_path}: {e}")
        return ""

def clean_text_artifacts(text: str) -> str:
    """Helper to clean PDF/PPTX artifacts."""
    replacements = {
        "ﬁ": "fi", "ﬂ": "fl", "ﬀ": "ff", "ﬃ": "ffi", "ﬄ": "ffl",
        "“": '"', "”": '"', "‘": "'", "’": "'",
        "–": "-", "—": "-",
        "\u200b": ""
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text

def validate_file(filename: str) -> bool:
    """Check allowed extensions."""
    return filename.lower().endswith(('.pdf', '.pptx', '.ppt'))
