"""
Text Processing Pipeline for Similarity Checking
Implements: PDF/PPTX extraction, text cleaning, chunking, and document metadata.

Usage:
    from text_pipeline import extract_text_from_file, clean_text, chunk_text, DocumentMetadata, process_document
"""
import os
import re
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

# PDF extraction
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

# PPTX extraction
try:
    from pptx import Presentation
except ImportError:
    Presentation = None


# =============================================================================
# METADATA STRUCTURE
# =============================================================================

@dataclass
class DocumentMetadata:
    """Metadata for a processed document (repository/organization)."""
    document_id: str
    file_name: str
    file_path: str
    num_chunks: int
    indexing_time: float  # seconds
    file_type: str  # 'pdf' or 'pptx'
    num_pages_or_slides: int = 0
    raw_text_length: int = 0

    def to_dict(self) -> dict:
        return {
            'document_id': self.document_id,
            'file_name': self.file_name,
            'file_path': self.file_path,
            'num_chunks': self.num_chunks,
            'indexing_time': self.indexing_time,
            'file_type': self.file_type,
            'num_pages_or_slides': self.num_pages_or_slides,
            'raw_text_length': self.raw_text_length,
        }


# =============================================================================
# PDF TEXT EXTRACTION
# =============================================================================

def extract_text_from_pdf_pypdf2(pdf_path: str) -> List[Tuple[int, str]]:
    """Extract text page by page using PyPDF2. Returns list of (page_number_1based, page_text)."""
    if PyPDF2 is None:
        raise ImportError("PyPDF2 is required. Install with: pip install PyPDF2")
    result = []
    with open(pdf_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
                result.append((i + 1, text))
            except Exception as e:
                result.append((i + 1, f"[Error extracting page: {e}]"))
    return result


def extract_text_from_pdf_pdfplumber(pdf_path: str) -> List[Tuple[int, str]]:
    """Extract text page by page using pdfplumber. Returns list of (page_number_1based, page_text)."""
    if pdfplumber is None:
        raise ImportError("pdfplumber is required. Install with: pip install pdfplumber")
    result = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            try:
                text = page.extract_text() or ""
                result.append((i + 1, text))
            except Exception as e:
                result.append((i + 1, f"[Error extracting page: {e}]"))
    return result


def extract_text_from_pdf(pdf_path: str, method: str = "pdfplumber") -> List[Tuple[int, str]]:
    """Extract text from PDF page by page. method: 'pypdf2' or 'pdfplumber'."""
    if method == "pypdf2" and PyPDF2:
        return extract_text_from_pdf_pypdf2(pdf_path)
    return extract_text_from_pdf_pdfplumber(pdf_path)


def pdf_pages_to_full_text(pages: List[Tuple[int, str]]) -> str:
    """Convert list of (page_num, text) to a single string with page breaks."""
    return "\n\n".join(text for _, text in pages)


# =============================================================================
# PPTX TEXT EXTRACTION
# =============================================================================

def extract_text_from_pptx(pptx_path: str) -> List[Tuple[int, str]]:
    """Extract text from PPTX slide by slide. Returns list of (slide_number_1based, slide_text)."""
    if Presentation is None:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    result = []
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        result.append((i + 1, "\n".join(parts)))
    return result


def pptx_slides_to_full_text(slides: List[Tuple[int, str]]) -> str:
    """Convert list of (slide_num, text) to a single string."""
    return "\n\n".join(text for _, text in slides)


# =============================================================================
# UNIFIED FILE EXTRACTION
# =============================================================================

def extract_text_from_file(
    file_path: str,
    pdf_method: str = "pdfplumber"
) -> Tuple[str, str, int]:
    """Extract text from PDF or PPTX. Returns (full_text, file_type, num_pages_or_slides)."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    ext = path.suffix.lower()
    if ext == ".pdf":
        pages = extract_text_from_pdf(file_path, method=pdf_method)
        full_text = pdf_pages_to_full_text(pages)
        return full_text, "pdf", len(pages)
    elif ext in (".pptx", ".ppt"):
        if ext == ".ppt":
            raise ValueError("Only .pptx is supported; .ppt (old format) is not supported.")
        slides = extract_text_from_pptx(file_path)
        full_text = pptx_slides_to_full_text(slides)
        return full_text, "pptx", len(slides)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Use .pdf or .pptx.")


# =============================================================================
# FILE PROCESSING LAYER (A): Header/Footer, Duplicate, Lowercase, Stopwords
# =============================================================================

def remove_duplicate_lines(text: str) -> str:
    """Remove consecutive duplicate lines (keeps one copy)."""
    if not text or not isinstance(text, str):
        return ""
    lines = text.split("\n")
    out = []
    prev = None
    for line in lines:
        s = line.strip()
        if s != prev:
            out.append(line)
            prev = s
    return "\n".join(out)


def remove_header_footer_lines(text: str) -> str:
    """Remove lines that look like page numbers or very short header/footer (e.g. only digits, or 1-2 words)."""
    if not text or not isinstance(text, str):
        return ""
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    filtered = []
    for ln in lines:
        if re.match(r"^\d+$", ln):
            continue
        if len(ln) <= 2:
            continue
        words = ln.split()
        if len(words) <= 1 and len(ln) < 20:
            continue
        filtered.append(ln)
    return "\n".join(filtered)


def normalize_lowercase(text: str) -> str:
    """Lowercase for consistent comparison (embedding model may be case-insensitive; helps lexical)."""
    return text.lower() if text else ""


# Common boilerplate phrases in weekly reports/templates that should not drive similarity.
_BOILERPLATE_PHRASES = [
    "weekly report",
    "weeklyreport",
    "tasks completed",
    "tasks com",
    "challenges faced",
    "progress made",
    "work narrative of the current week",
    "follow-up from last week",
    "follow up from last week",
]


def remove_boilerplate_phrases(text: str) -> str:
    """Remove known boilerplate phrases (case-insensitive) to reduce false matches."""
    if not text or not isinstance(text, str):
        return text or ""
    t = text
    for phrase in _BOILERPLATE_PHRASES:
        t = re.sub(re.escape(phrase), " ", t, flags=re.IGNORECASE)
    # Collapse multiple spaces/newlines introduced by removals.
    t = re.sub(r"\s+", " ", t)
    return t.strip()


# Patterns and phrases for institutional metadata (student info, course codes, etc.).
_METADATA_IGNORE_PATTERNS = [
    # Student identifiers
    r"\b(?:student|id|identification|roll)\s*(?:no|number|#)?\s*:?\s*\d+",
    r"\b(?:student|candidate)\s+(?:name|id)\s*:",
    # Dates
    r"\b(?:submission|submitted|date)\s*:?\s*\d{1,2}[/-]\d{1,2}[/-]\d{2,4}",
    r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b",
    r"\b(?:january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{1,2},?\s+\d{4}",
    # Contact information
    r"\bemail\s*:?\s*\S+@\S+\.\S+",
    r"\b(?:phone|mobile)\s*:?\s*[\d\s\-\(\)\+]+",
    # Course codes / academic session
    r"\bcourse\s+(?:code|number|id)\s*:?\s*[A-Z]{2,4}\s*\d{3,4}",
    r"\bsemester\s*:?\s*\w+",
    r"\bsession\s*:?\s*\d{4}\s*-\s*\d{4}",
    # Registration / enrollment
    r"\b(?:registration|enrollment)\s+(?:number|no|#)\s*:?\s*\d+",
]

_METADATA_IGNORE_PHRASES = {
    "submitted by",
    "submitted to",
    "prepared by",
    "prepared for",
    "course code",
    "course name",
    "course title",
    "instructor",
    "professor",
    "supervisor",
    "advisor",
    "department of",
    "faculty of",
    "university",
    "college",
    "institute",
    "academic year",
    "registration number",
    "enrollment number",
    "student id",
    "in partial fulfillment",
    "partial fulfillment of",
    "requirement for the degree",
}


def remove_institutional_metadata(text: str) -> str:
    """Remove lines that look like institutional metadata to avoid inflating similarity."""
    if not text or not isinstance(text, str):
        return text or ""
    lines = text.split("\n")
    kept = []
    for ln in lines:
        stripped = ln.strip()
        if not stripped:
            continue
        lower = stripped.lower()
        # Phrase-based ignore
        if any(phrase in lower for phrase in _METADATA_IGNORE_PHRASES):
            continue
        # Regex-based ignore
        ignore = False
        for pattern in _METADATA_IGNORE_PATTERNS:
            if re.search(pattern, stripped, flags=re.IGNORECASE):
                ignore = True
                break
        if ignore:
            continue
        kept.append(stripped)
    return "\n".join(kept)


def remove_references_section(text: str) -> str:
    """Trim off the references/bibliography section so common citation lists don't dominate matches."""
    if not text or not isinstance(text, str):
        return text or ""
    ref_headers = [
        r"\bReferences\b",
        r"\bBibliography\b",
        r"\bWorks Cited\b",
        r"\bLiterature Cited\b",
        r"\bCitations\b",
    ]
    earliest_pos = len(text)
    for header in ref_headers:
        m = re.search(header, text, flags=re.IGNORECASE)
        if m:
            earliest_pos = min(earliest_pos, m.start())
    if earliest_pos < len(text):
        return text[:earliest_pos]
    return text


# Optional stopwords (common English); set to None to disable.
_STOPWORDS = None

def _get_stopwords() -> set:
    global _STOPWORDS
    if _STOPWORDS is None:
        try:
            import nltk
            try:
                nltk.data.find("corpora/stopwords")
            except LookupError:
                nltk.download("stopwords", quiet=True)
            from nltk.corpus import stopwords
            _STOPWORDS = set(stopwords.words("english"))
        except Exception:
            _STOPWORDS = set()
    return _STOPWORDS


def remove_stopwords(text: str, stopwords_optional: bool = True) -> str:
    """Remove stopwords (optional). When enabled, reduces noise for similarity."""
    if not stopwords_optional or not text:
        return text or ""
    sw = _get_stopwords()
    if not sw:
        return text
    words = text.split()
    kept = [w for w in words if w.lower() not in sw]
    return " ".join(kept)


def apply_file_processing_layer(
    full_text: str,
    lowercase: bool = True,
    remove_stopwords_opt: bool = False,
) -> str:
    """
    (A) File Processing Layer: duplicate line remover, header/footer remover,
    clean, lowercase normalize, optional stopword remove.
    """
    if not full_text or not full_text.strip():
        return ""
    t = remove_duplicate_lines(full_text)
    t = remove_header_footer_lines(t)
    t = remove_references_section(t)
    t = remove_institutional_metadata(t)
    # Remove boilerplate phrases common in weekly reports / templates so they do not dominate similarity.
    t = remove_boilerplate_phrases(t)
    t = clean_text(t)
    if lowercase:
        t = normalize_lowercase(t)
    if remove_stopwords_opt:
        t = remove_stopwords(t, stopwords_optional=True)
    return t.strip()


# =============================================================================
# TEXT CLEANING
# =============================================================================

def clean_text(text: str) -> str:
    """Clean extracted text: normalize whitespace, remove control characters."""
    if not text or not isinstance(text, str):
        return ""
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    return text.strip()


def clean_text_aggressive(text: str, keep_punctuation: bool = True) -> str:
    """More aggressive cleaning: keep only letters, digits, and optionally punctuation."""
    text = clean_text(text)
    if keep_punctuation:
        text = re.sub(r"[^\w\s.,!?;:\-'\"()]", "", text)
    else:
        text = re.sub(r"[^\w\s]", "", text)
    return re.sub(r"\s+", " ", text).strip()


# =============================================================================
# CHUNKING STRATEGY
# =============================================================================

def chunk_by_words(
    text: str,
    max_words: int = 100,
    overlap_words: int = 0,
    min_chunk_words: int = 5
) -> List[str]:
    """Split text into chunks by word count with optional overlap."""
    text = clean_text(text)
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + max_words, len(words))
        chunk_words = words[start:end]
        if len(chunk_words) >= min_chunk_words or end == len(words):
            chunks.append(" ".join(chunk_words))
        if overlap_words <= 0 or end >= len(words):
            start = end
        else:
            start = end - overlap_words
    return chunks


def chunk_by_sentences(text: str, min_length: int = 3) -> List[str]:
    """Split by sentence boundaries (. ? !). Each chunk = one sentence."""
    text = clean_text(text)
    if not text:
        return []
    parts = re.split(r"(?<=[.!?])\s*(?=[A-Z]|$)", text)
    sentences = [s.strip() for s in parts if s and len(s.strip()) >= min_length]
    if len(sentences) <= 1 and "\n" in text:
        parts = [p.strip() for p in text.split("\n") if p.strip() and len(p.strip()) >= min_length]
        if parts:
            return parts
    return sentences if sentences else [text.strip()] if text.strip() else []


def chunk_by_paragraphs(full_text: str, fallback_max_words: int = 80) -> List[str]:
    """
    Split by paragraph (double newline). Each chunk = one paragraph so comparison is meaningful.
    If no paragraph breaks, fall back to small word chunks (80 words).
    """
    if not full_text or not full_text.strip():
        return []
    # Split on double newline (one or more newlines with optional spaces) before we collapse whitespace.
    raw_paragraphs = re.split(r"\n\s*\n", full_text)
    cleaned = [clean_text(p).strip() for p in raw_paragraphs if p and clean_text(p).strip()]
    if len(cleaned) >= 2:
        return cleaned
    # No paragraph breaks: use smaller word chunks so "Your text" / "Matched" are comparable.
    single = clean_text(full_text)
    return chunk_by_words(single, max_words=fallback_max_words, overlap_words=0, min_chunk_words=10)


def chunk_by_tokens_approx(
    text: str,
    max_tokens_approx: int = 256,
    overlap_tokens_approx: int = 0,
    min_chunk_tokens_approx: int = 10
) -> List[str]:
    """Chunk by approximate token count (~4 chars per token)."""
    approx_chars_per_token = 4
    max_chars = max_tokens_approx * approx_chars_per_token
    overlap_chars = overlap_tokens_approx * approx_chars_per_token
    min_chars = min_chunk_tokens_approx * approx_chars_per_token
    text = clean_text(text)
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        segment = text[start:end]
        if end < len(text):
            last_space = segment.rfind(" ")
            if last_space > max_chars // 2:
                end = start + last_space + 1
                segment = text[start:end]
        if len(segment.strip()) >= min_chars or end >= len(text):
            chunks.append(segment.strip())
        if overlap_chars <= 0 or end >= len(text):
            start = end
        else:
            start = end - overlap_chars
    return [c for c in chunks if c]


# =============================================================================
# FULL PIPELINE: EXTRACT -> CLEAN -> CHUNK -> METADATA
# =============================================================================

def process_document(
    file_path: str,
    pdf_method: str = "pdfplumber",
    chunk_strategy: str = "words",
    max_chunk_size: int = 200,
    overlap: int = 20,
    document_id: Optional[str] = None,
) -> Tuple[List[str], DocumentMetadata, str]:
    """
    Full pipeline: extract, clean, chunk, build metadata.
    Returns (chunks, DocumentMetadata, cleaned_full_text).
    """
    start_time = time.time()
    doc_id = document_id or str(uuid.uuid4())[:8]
    file_path = os.path.abspath(file_path)
    file_name = os.path.basename(file_path)

    full_text, file_type, num_pages = extract_text_from_file(file_path, pdf_method=pdf_method)
    raw_length = len(full_text)
    # (A) File processing layer: duplicate lines, header/footer, clean, lowercase, optional stopwords
    cleaned = apply_file_processing_layer(full_text, lowercase=True, remove_stopwords_opt=False)

    # (B) Chunking: 200-300 word chunks, 20-30% overlap (do not embed whole document)
    if chunk_strategy == "paragraphs":
        chunks = chunk_by_paragraphs(full_text)
    elif chunk_strategy == "sentences":
        chunks = chunk_by_sentences(cleaned)
    elif chunk_strategy == "tokens_approx":
        chunks = chunk_by_tokens_approx(
            cleaned,
            max_tokens_approx=max_chunk_size,
            overlap_tokens_approx=overlap
        )
    else:
        chunks = chunk_by_words(
            cleaned,
            max_words=max_chunk_size,
            overlap_words=overlap
        )

    indexing_time = round(time.time() - start_time, 4)
    meta = DocumentMetadata(
        document_id=doc_id,
        file_name=file_name,
        file_path=file_path,
        num_chunks=len(chunks),
        indexing_time=indexing_time,
        file_type=file_type,
        num_pages_or_slides=num_pages,
        raw_text_length=raw_length,
    )
    return chunks, meta, cleaned
